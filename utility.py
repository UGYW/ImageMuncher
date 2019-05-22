# UTILIY FUNCTIONS FOR GETTING IMAGE DATA
from constants import *
import os
from pptx.util import Cm, Pt
import re
import os

# Reference: https://python-pptx.readthedocs.io/en/latest/api/shapes.html

def get_sub_folder_paths(path_to_main_folder):
    contents = os.listdir(path_to_main_folder)
    # Adds the path to the main folder in front for traversal
    contents = [path_to_main_folder + "/" + i for i in contents if bool(re.match('[\d/-_]+$', i))]
    return contents

def get_sub_folder_content_paths(path_to_sub_folder):
    # Make a dictionary based on type
    res = {}
    contents = os.listdir(path_to_sub_folder)
    res[RGB] = [path_to_sub_folder + "/" + i for i in contents if RGB in i][0]
    res[OXY] = [path_to_sub_folder + "/" + i for i in contents if OXY in i][0]
    res[THI] = [path_to_sub_folder + "/" + i for i in contents if THI in i][0]
    res[NIR] = [path_to_sub_folder + "/" + i for i in contents if NIR in i][0]
    res[TWI] = [path_to_sub_folder + "/" + i for i in contents if TWI in i][0]
    return res

def get_caption_text(path_to_main_folder):
    # Gets what is assumed to be the only txt file in the main folder
    contents = os.listdir(path_to_main_folder)
    path_to_text = [path_to_main_folder + "/" + i for i in contents if ".txt" in i][0]
    with open(path_to_text, 'r') as myfile:
        text = myfile.read()
    return text

def get_date(path):
    return os.path.basename(os.path.normpath(path))

def add_picture(slide, picture, location, size = (-1, -1)):
    width = None
    height = None
    if size[0] > 0: width = size[0]
    if size[1] > 0: height = size[1]
    slide.shapes.add_picture(picture, location[0], location[1], width, height)

def add_title(slide, title_text):
    slide.shapes.title.text = title_text
    slide.shapes.title.font = Pt(12)

FONT_SIZE = 12
def add_text(slide, text, location, size = (-1, -1)):
    width = None
    height = None
    if size[0] > 0: width = size[0]
    if size[1] > 0: height = size[1]
    textbox = slide.shapes.add_textbox(location[0], location[1], width, height)
    textbox.text = text
    textbox.size = Pt(12)

# https://github.com/scanny/python-pptx/issues/195
def add_background(prs, slide, background_pic):
    left = top = Cm(0)
    pic = slide.shapes.add_picture(background_pic, left, top,
                                    width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

def add_table(slide, rows, cols, location, size = (-1, -1)):
    width = None
    height = None
    if size[0] > 0: width = size[0]
    if size[1] > 0: height = size[1]
    slide.shapes.add_table(rows, cols, location[0], location[1], width, height)