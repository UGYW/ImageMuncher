from pptx import *
from INPUTS import *
from utility import *
import os

def main():
    # Make a new presentation
    prs = Presentation()

    sub_folder_paths = get_sub_folder_paths(PATH_TO_DIRECTORIES)

    caption_text = get_caption_text(PATH_TO_DIRECTORIES)

    # Make a new slide for each sub folder and its contents
    for sub_folder_path in sub_folder_paths:
        sub_folder_content_paths = get_sub_folder_content_paths(sub_folder_path)

        # Make a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_CONTENT])  # title and content layout

        # Add the background image
        add_background(prs, slide, PATH_TO_BACKGROUND_PIC)

        # Add the title - using the sub folder path
        title_text = get_date(sub_folder_path)
        add_title(slide, title_text)

        # Add the caption text
        add_text(slide, caption_text, CAPTION_TXT_LOC, CAPTION_IMAGE_SIZE)

        # Add the table
        add_table(slide, TABLE_ROW_NUM, TABLE_COL_NUM, TABLE_LOC, TABLE_SIZE)

        # Add the caption (the RGB image) image
        add_picture(slide, sub_folder_content_paths[RGB], CAPTION_IMAGE_LOC, CAPTION_IMAGE_SIZE)

        # Add the four images
        add_picture(slide, sub_folder_content_paths[OXY], UPPER_LEFT, MAIN_IMAGE_SIZE)
        add_picture(slide, sub_folder_content_paths[THI], UPPER_RIGHT, MAIN_IMAGE_SIZE)
        add_picture(slide, sub_folder_content_paths[NIR], LOWER_LEFT, MAIN_IMAGE_SIZE)
        add_picture(slide, sub_folder_content_paths[TWI], LOWER_RIGHT, MAIN_IMAGE_SIZE)

    prs.save(PATH_TO_POWERPOINT)

if __name__ == '__main__':
    main()